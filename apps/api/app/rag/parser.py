"""Document parsing — PDF/DOCX/PPTX/TXT/CSV/HTML/URL/manual (spec §2.1, §11.2, §11.3).

Backends sit behind `ParserBackend`, so the heavy, native-dependency parsers can live
in the worker image (AMD AUP side, §72) while the API process only ever holds the
protocol. Every backend returns `ParsedDocument` — a flat block list plus page count,
metadata and the character count used by the OCR gate.

Only the text/CSV/HTML/URL backends are dependency-free. `pypdf`, `python-docx` and
`python-pptx` are imported lazily and their absence produces a clear
`ParserUnavailableError` rather than an import-time crash, so a deployment that has
not installed them still runs (and the document lands in `failed` with a readable
`failure_reason`, per §29).
"""

from __future__ import annotations

import csv
import html
import io
import re
from collections.abc import Sequence
from enum import StrEnum
from typing import Any, Protocol, runtime_checkable

import structlog
from pydantic import BaseModel, ConfigDict, Field

from app.rag.ocr import NullOcr, OcrPort, merge_ocr, should_ocr
from app.rag.structure import (
    Block,
    BlockKind,
    blocks_from_text,
    build_section_tree,
    is_table_well_formed,
    outline,
    table_to_text,
)

log = structlog.get_logger(__name__)

#: Mirrors `DocumentSourceKind` in packages/shared/src/entities.ts
class SourceKind(StrEnum):
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    TXT = "txt"
    CSV = "csv"
    HTML = "html"
    URL = "url"
    MANUAL = "manual"


MIME_BY_KIND: dict[SourceKind, tuple[str, ...]] = {
    SourceKind.PDF: ("application/pdf",),
    SourceKind.DOCX: (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ),
    SourceKind.PPTX: (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ),
    SourceKind.TXT: ("text/plain", "text/markdown"),
    SourceKind.CSV: ("text/csv", "application/csv"),
    SourceKind.HTML: ("text/html", "application/xhtml+xml"),
}

MAX_UPLOAD_BYTES = 100 * 1024 * 1024


class ParserError(RuntimeError):
    """Parsing failed for a reason worth showing the uploader (§29 failure_reason)."""


class ParserUnavailableError(ParserError):
    """The backend for this format is not installed in this image."""


class DocumentPayload(BaseModel):
    """What the worker hands the parser."""

    model_config = ConfigDict(extra="forbid")

    filename: str = ""
    source_kind: SourceKind = SourceKind.TXT
    content: bytes = b""
    #: for `manual` and `url` kinds
    text: str = ""
    url: str = ""
    mime_type: str | None = None
    encoding: str = "utf-8"
    language: str = "zh-TW"


class ParsedDocument(BaseModel):
    model_config = ConfigDict(extra="forbid")

    blocks: list[Block] = Field(default_factory=list)
    page_count: int = 0
    text_chars: int = 0
    parser: str = ""
    ocr_applied: bool = False
    parser_confidence: float = 1.0
    metadata: dict[str, Any] = Field(default_factory=dict)
    outline: list[dict[str, Any]] = Field(default_factory=list)

    @property
    def text(self) -> str:
        return "\n".join(b.text for b in self.blocks)


@runtime_checkable
class ParserBackend(Protocol):
    name: str

    def supports(self, kind: SourceKind) -> bool: ...

    async def parse(self, payload: DocumentPayload) -> ParsedDocument: ...


def validate(payload: DocumentPayload) -> None:
    """§11.2 checks that belong to the parser: MIME + size + emptiness.

    Virus scan, duplicate detection and the encryption policy live in
    `KnowledgeService.register_upload` because they need tenant configuration.
    """
    if payload.source_kind in (SourceKind.MANUAL,):
        if not payload.text.strip():
            raise ParserError("manual text is empty")
        return
    if payload.source_kind is SourceKind.URL:
        if not re.match(r"^https?://", payload.url):
            raise ParserError("url must be http(s)")
        return
    if not payload.content:
        raise ParserError("uploaded file is empty")
    if len(payload.content) > MAX_UPLOAD_BYTES:
        raise ParserError(f"file exceeds {MAX_UPLOAD_BYTES // (1024 * 1024)}MB limit")
    expected = MIME_BY_KIND.get(payload.source_kind)
    if expected and payload.mime_type and payload.mime_type not in expected:
        raise ParserError(
            f"mime type {payload.mime_type} does not match source kind {payload.source_kind}"
        )


# ---------------------------------------------------------------------------
# dependency-free backends
# ---------------------------------------------------------------------------
class PlainTextBackend:
    name = "plaintext"

    def supports(self, kind: SourceKind) -> bool:
        return kind in (SourceKind.TXT, SourceKind.MANUAL)

    async def parse(self, payload: DocumentPayload) -> ParsedDocument:
        text = payload.text or _decode(payload.content, payload.encoding)
        blocks = blocks_from_text(text, page=1)
        return _assemble(blocks, page_count=1, parser=self.name)


class CsvBackend:
    """CSV becomes one table block per 200 rows so large sheets stay chunkable."""

    name = "csv"
    rows_per_block = 200

    def supports(self, kind: SourceKind) -> bool:
        return kind is SourceKind.CSV

    async def parse(self, payload: DocumentPayload) -> ParsedDocument:
        text = payload.text or _decode(payload.content, payload.encoding)
        sample = text[:4096]
        try:
            dialect: Any = csv.Sniffer().sniff(sample, delimiters=",;\t|")
        except csv.Error:
            dialect = csv.excel
        rows = [row for row in csv.reader(io.StringIO(text), dialect) if any(c.strip() for c in row)]
        if not rows:
            raise ParserError("csv contained no data rows")
        header, body = rows[0], rows[1:]
        blocks: list[Block] = []
        for start in range(0, max(len(body), 1), self.rows_per_block):
            window = [header, *body[start : start + self.rows_per_block]]
            blocks.append(
                Block(
                    kind=BlockKind.TABLE,
                    text=table_to_text(window, caption=payload.filename),
                    rows=[list(r) for r in window],
                    page=1,
                    order=len(blocks),
                    section_path=[payload.filename] if payload.filename else [],
                )
            )
        return _assemble(
            blocks,
            page_count=1,
            parser=self.name,
            metadata={"columns": header, "row_count": len(body)},
        )


_SCRIPT_STYLE = re.compile(r"<(script|style|noscript)[^>]*>.*?</\1>", re.IGNORECASE | re.DOTALL)
_TAG = re.compile(r"<[^>]+>")
_BLOCK_TAG = re.compile(
    r"</?(p|div|section|article|br|li|tr|h[1-6]|table|thead|tbody|blockquote)[^>]*>",
    re.IGNORECASE,
)
_HEADING_TAG = re.compile(r"<h([1-6])[^>]*>(.*?)</h\1>", re.IGNORECASE | re.DOTALL)
_TABLE_TAG = re.compile(r"<table[^>]*>(.*?)</table>", re.IGNORECASE | re.DOTALL)
_ROW_TAG = re.compile(r"<tr[^>]*>(.*?)</tr>", re.IGNORECASE | re.DOTALL)
_CELL_TAG = re.compile(r"<(t[dh])[^>]*>(.*?)</\1>", re.IGNORECASE | re.DOTALL)
_TITLE_TAG = re.compile(r"<title[^>]*>(.*?)</title>", re.IGNORECASE | re.DOTALL)


class HtmlBackend:
    """Regex-based HTML extraction — deliberately dependency-free.

    Enterprise sources here are FAQ pages and intranet SOPs, not arbitrary SPA sites,
    so headings/tables/paragraphs are all we need. Swap in a `SelectolaxBackend` later
    without touching the pipeline: it only depends on the `ParserBackend` protocol.
    """

    name = "html"

    def supports(self, kind: SourceKind) -> bool:
        return kind in (SourceKind.HTML, SourceKind.URL)

    async def parse(self, payload: DocumentPayload) -> ParsedDocument:
        raw = payload.text or _decode(payload.content, payload.encoding)
        title_match = _TITLE_TAG.search(raw)
        body = _SCRIPT_STYLE.sub(" ", raw)

        tables: list[list[list[str]]] = []
        for table_html in _TABLE_TAG.findall(body):
            rows: list[list[str]] = []
            for row_html in _ROW_TAG.findall(table_html):
                cells = [_clean(cell) for _tag, cell in _CELL_TAG.findall(row_html)]
                if any(cells):
                    rows.append(cells)
            if is_table_well_formed(rows):
                tables.append(rows)
        body_no_tables = _TABLE_TAG.sub("\n", body)

        # Convert headings to markdown so `blocks_from_text` can level them.
        def heading_repl(match: re.Match[str]) -> str:
            level = int(match.group(1))
            return f"\n{'#' * level} {_clean(match.group(2))}\n"

        marked = _HEADING_TAG.sub(heading_repl, body_no_tables)
        marked = _BLOCK_TAG.sub("\n", marked)
        text = _clean_multiline(_TAG.sub(" ", marked))
        blocks = blocks_from_text(text, page=1)
        for rows in tables:
            blocks.append(
                Block(
                    kind=BlockKind.TABLE,
                    text=table_to_text(rows),
                    rows=rows,
                    page=1,
                    order=len(blocks),
                )
            )
        return _assemble(
            blocks,
            page_count=1,
            parser=self.name,
            metadata={
                "title": _clean(title_match.group(1)) if title_match else "",
                "source_url": payload.url,
            },
        )


class UrlBackend:
    """Fetch a URL then delegate to the HTML backend (§11.2 Import URL)."""

    name = "url"

    def __init__(self, *, timeout_s: float = 20.0, max_bytes: int = 8 * 1024 * 1024) -> None:
        self._timeout = timeout_s
        self._max_bytes = max_bytes
        self._html = HtmlBackend()

    def supports(self, kind: SourceKind) -> bool:
        return kind is SourceKind.URL

    async def parse(self, payload: DocumentPayload) -> ParsedDocument:
        import httpx

        async with httpx.AsyncClient(timeout=self._timeout, follow_redirects=True) as client:
            response = await client.get(payload.url, headers={"User-Agent": "AICoach-KB/1.0"})
        if response.status_code >= 400:
            raise ParserError(f"url returned {response.status_code}")
        content = response.content[: self._max_bytes]
        parsed = await self._html.parse(
            payload.model_copy(update={"content": content, "text": "", "url": payload.url})
        )
        parsed.metadata["source_url"] = payload.url
        parsed.metadata["fetched_status"] = response.status_code
        parsed.parser = self.name
        return parsed


# ---------------------------------------------------------------------------
# optional native backends
# ---------------------------------------------------------------------------
class PdfBackend:
    """`pypdf` text extraction, page by page (page numbers feed §12.5 citations)."""

    name = "pdf"

    def supports(self, kind: SourceKind) -> bool:
        return kind is SourceKind.PDF

    async def parse(self, payload: DocumentPayload) -> ParsedDocument:
        import asyncio

        return await asyncio.to_thread(self._parse_sync, payload)

    def _parse_sync(self, payload: DocumentPayload) -> ParsedDocument:
        try:
            from pypdf import PdfReader  # type: ignore[import-not-found]
        except ImportError as exc:  # pragma: no cover - optional dependency
            raise ParserUnavailableError(
                "pypdf is not installed in this image; PDF ingestion is unavailable"
            ) from exc
        reader = PdfReader(io.BytesIO(payload.content))
        blocks: list[Block] = []
        pages: dict[int, str] = {}
        for index, page in enumerate(reader.pages, start=1):
            try:
                text = page.extract_text() or ""
            except Exception as exc:  # noqa: BLE001 - one bad page must not kill the doc
                log.warning("pdf.page_failed", page=index, error=repr(exc))
                text = ""
            pages[index] = text
            blocks.extend(blocks_from_text(text, page=index, start_order=len(blocks)))
        meta = {}
        try:
            info = reader.metadata or {}
            meta = {str(k).lstrip("/"): str(v) for k, v in dict(info).items()}
        except Exception:  # noqa: BLE001 - metadata is best-effort
            meta = {}
        parsed = _assemble(blocks, page_count=len(pages), parser=self.name, metadata=meta)
        parsed.metadata["pages_text"] = {str(k): len(v) for k, v in pages.items()}
        return parsed


class DocxBackend:
    name = "docx"

    def supports(self, kind: SourceKind) -> bool:
        return kind is SourceKind.DOCX

    async def parse(self, payload: DocumentPayload) -> ParsedDocument:
        import asyncio

        return await asyncio.to_thread(self._parse_sync, payload)

    def _parse_sync(self, payload: DocumentPayload) -> ParsedDocument:
        try:
            import docx  # type: ignore[import-not-found]
        except ImportError as exc:  # pragma: no cover
            raise ParserUnavailableError(
                "python-docx is not installed in this image; DOCX ingestion is unavailable"
            ) from exc
        document = docx.Document(io.BytesIO(payload.content))
        blocks: list[Block] = []
        path: list[str] = []
        for paragraph in document.paragraphs:
            text = (paragraph.text or "").strip()
            if not text:
                continue
            style = (paragraph.style.name if paragraph.style is not None else "") or ""
            level = _docx_heading_level(style)
            if level:
                del path[level - 1 :]
                path.append(text)
                blocks.append(
                    Block(
                        kind=BlockKind.HEADING,
                        text=text,
                        level=level,
                        order=len(blocks),
                        section_path=list(path),
                    )
                )
            elif style.startswith("List"):
                blocks.append(
                    Block(
                        kind=BlockKind.LIST_ITEM,
                        text=text,
                        order=len(blocks),
                        section_path=list(path),
                    )
                )
            else:
                blocks.append(
                    Block(
                        kind=BlockKind.PARAGRAPH,
                        text=text,
                        order=len(blocks),
                        section_path=list(path),
                    )
                )
        for table in document.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if rows:
                blocks.append(
                    Block(
                        kind=BlockKind.TABLE,
                        text=table_to_text(rows),
                        rows=rows,
                        order=len(blocks),
                        section_path=list(path),
                    )
                )
        return _assemble(blocks, page_count=1, parser=self.name)


class PptxBackend:
    """Slides: title becomes a heading, body text becomes paragraphs, one page each."""

    name = "pptx"

    def supports(self, kind: SourceKind) -> bool:
        return kind is SourceKind.PPTX

    async def parse(self, payload: DocumentPayload) -> ParsedDocument:
        import asyncio

        return await asyncio.to_thread(self._parse_sync, payload)

    def _parse_sync(self, payload: DocumentPayload) -> ParsedDocument:
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
        except ImportError as exc:  # pragma: no cover
            raise ParserUnavailableError(
                "python-pptx is not installed in this image; PPTX ingestion is unavailable"
            ) from exc
        presentation = Presentation(io.BytesIO(payload.content))
        blocks: list[Block] = []
        for index, slide in enumerate(presentation.slides, start=1):
            title = ""
            if slide.shapes.title is not None:
                title = (slide.shapes.title.text or "").strip()
            if title:
                blocks.append(
                    Block(
                        kind=BlockKind.HEADING,
                        text=title,
                        level=1,
                        page=index,
                        order=len(blocks),
                        section_path=[title],
                    )
                )
            for shape in slide.shapes:
                if shape == slide.shapes.title or not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        blocks.append(
                            Block(
                                kind=BlockKind.PARAGRAPH,
                                text=text,
                                page=index,
                                order=len(blocks),
                                section_path=[title] if title else [],
                            )
                        )
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                blocks.append(
                    Block(
                        kind=BlockKind.CAPTION,
                        text=notes,
                        page=index,
                        order=len(blocks),
                        section_path=[title] if title else [],
                        metadata={"speaker_notes": True},
                    )
                )
        return _assemble(
            blocks, page_count=len(presentation.slides), parser=self.name
        )


# ---------------------------------------------------------------------------
# registry
# ---------------------------------------------------------------------------
class DocumentParser:
    """Backend registry + the §11.3 extract/OCR sequence."""

    def __init__(
        self,
        backends: Sequence[ParserBackend] | None = None,
        *,
        ocr: OcrPort | None = None,
    ) -> None:
        self.backends: list[ParserBackend] = list(
            backends
            or (
                PlainTextBackend(),
                CsvBackend(),
                HtmlBackend(),
                UrlBackend(),
                PdfBackend(),
                DocxBackend(),
                PptxBackend(),
            )
        )
        self.ocr: OcrPort = ocr or NullOcr()

    def backend_for(self, kind: SourceKind) -> ParserBackend:
        for backend in self.backends:
            if backend.supports(kind):
                return backend
        raise ParserUnavailableError(f"no parser backend for source kind {kind}")

    async def parse(self, payload: DocumentPayload) -> ParsedDocument:
        validate(payload)
        backend = self.backend_for(payload.source_kind)
        parsed = await backend.parse(payload)

        if should_ocr(
            text_chars=parsed.text_chars,
            page_count=parsed.page_count,
            parser_confidence=parsed.parser_confidence,
            source_kind=str(payload.source_kind),
        ):
            log.info(
                "parser.ocr_triggered",
                filename=payload.filename,
                text_chars=parsed.text_chars,
                pages=parsed.page_count,
            )
            ocr_result = await self.ocr.recognise(
                payload.content, languages=language_hint(payload)
            )
            if ocr_result.pages:
                merged = merge_ocr(ocr_result)
                blocks: list[Block] = list(parsed.blocks)
                for page_number, text in sorted(merged.items()):
                    blocks.extend(
                        blocks_from_text(
                            text,
                            page=page_number,
                            start_order=len(blocks),
                            confidence=max(0.1, ocr_result.mean_confidence),
                        )
                    )
                parsed = _assemble(
                    blocks,
                    page_count=max(parsed.page_count, len(merged)),
                    parser=f"{parsed.parser}+{ocr_result.engine}",
                    metadata=parsed.metadata,
                )
                parsed.ocr_applied = True
                parsed.parser_confidence = max(0.1, ocr_result.mean_confidence)
        if not parsed.blocks:
            raise ParserError("no extractable text (document may be an unreadable scan)")
        return parsed


def language_hint(payload: DocumentPayload) -> tuple[str, ...]:
    """Tesseract language codes for the document's declared locale."""
    return ("chi_tra", "eng") if payload.language.startswith("zh") else ("eng",)


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------
def _assemble(
    blocks: Sequence[Block],
    *,
    page_count: int,
    parser: str,
    metadata: dict[str, Any] | None = None,
) -> ParsedDocument:
    ordered = [b.model_copy(update={"order": i}) for i, b in enumerate(blocks)]
    tree = build_section_tree(ordered)
    return ParsedDocument(
        blocks=ordered,
        page_count=page_count,
        text_chars=sum(len(b.text) for b in ordered),
        parser=parser,
        metadata=dict(metadata or {}),
        outline=outline(tree),
    )


def _decode(content: bytes, encoding: str) -> str:
    for candidate in (encoding, "utf-8", "utf-8-sig", "big5", "cp950", "gb18030", "latin-1"):
        try:
            return content.decode(candidate)
        except (UnicodeDecodeError, LookupError):
            continue
    return content.decode("utf-8", errors="replace")


def _clean(fragment: str) -> str:
    return re.sub(r"\s+", " ", html.unescape(_TAG.sub(" ", fragment))).strip()


def _clean_multiline(text: str) -> str:
    unescaped = html.unescape(text)
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in unescaped.split("\n")]
    out: list[str] = []
    for line in lines:
        if not line and out and not out[-1]:
            continue
        out.append(line)
    return "\n".join(out)


def _docx_heading_level(style_name: str) -> int:
    match = re.match(r"Heading\s*(\d)", style_name, re.IGNORECASE)
    if match:
        return int(match.group(1))
    if style_name.lower() in ("title", "subtitle"):
        return 1
    return 0


__all__ = [
    "MAX_UPLOAD_BYTES",
    "MIME_BY_KIND",
    "CsvBackend",
    "DocumentParser",
    "DocumentPayload",
    "DocxBackend",
    "HtmlBackend",
    "ParsedDocument",
    "ParserBackend",
    "ParserError",
    "ParserUnavailableError",
    "PdfBackend",
    "PlainTextBackend",
    "PptxBackend",
    "SourceKind",
    "UrlBackend",
    "validate",
]
